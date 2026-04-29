import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Load data
with open("/home/benliangcs/tw-gov-video/output/selected_articles.json", "r", encoding="utf-8") as f:
    data = json.load(f)

# Create presentation
prs = Presentation()

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "GovClaw 自我介紹與工作匯報"
subtitle.text = "GCP 部署 x 跨平台通訊 x 新聞摘要影片自動化"

# Add slides for each section
bullet_slide_layout = prs.slide_layouts[1]

for item in data['selected']:
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = item['title']
    
    tf = body_shape.text_frame
    tf.text = item['script']
    
    # Optional styling
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(24)

# Save
output_path = "/home/benliangcs/tw-gov-video/output/GovClaw_Intro.pptx"
prs.save(output_path)
print(f"PPTX saved to {output_path}")
